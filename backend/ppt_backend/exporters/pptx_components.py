from __future__ import annotations

from dataclasses import dataclass
from typing import Any, Dict, List, Protocol, Tuple

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..domain.render_tree import ComponentSpec
from ..domain.theme import ThemeTokens
from ..registry.base import Registry


def _hex_to_rgb(hex_color: str) -> RGBColor:
    s = (hex_color or "").strip()
    if s.startswith("#"):
        s = s[1:]
    if len(s) == 3:
        s = "".join([c * 2 for c in s])
    if len(s) != 6:
        return RGBColor(255, 255, 255)
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _align(align: str):
    if align == "center":
        return PP_ALIGN.CENTER
    if align == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


@dataclass(frozen=True)
class PptxRenderContext:
    theme: ThemeTokens
    slide_px: Tuple[int, int] = (1280, 720)
    slide_inches: Tuple[float, float] = (13.333, 7.5)

    def to_inches(self, x_px: float, y_px: float, w_px: float, h_px: float):
        sw_px, sh_px = self.slide_px
        sw_in, sh_in = self.slide_inches
        sx = sw_in / sw_px
        sy = sh_in / sh_px
        return Inches(x_px * sx), Inches(y_px * sy), Inches(w_px * sx), Inches(h_px * sy)


class ComponentRenderer(Protocol):
    type: str

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None: ...


def build_component_renderer_registry() -> Registry[ComponentRenderer]:
    reg: Registry[ComponentRenderer] = Registry()
    for r in [
        TitleRenderer(),
        SubtitleRenderer(),
        TextRenderer(),
        BulletListRenderer(),
        QuoteRenderer(),
        DividerRenderer(),
        KpiCardsRenderer(),
        TimelineRenderer(),
        ComparisonTableRenderer(),
        SwotRenderer(),
        RoadmapRenderer(),
        ProcessFlowRenderer(),
        ChartRenderer(),
        MultiColumnRenderer(),
        TeamCardsRenderer(),
        ArchitectureDiagramRenderer(),
    ]:
        reg.register(r.type, r)
    return reg


class TitleRenderer:
    type = "Title"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        text = str(comp.props.get("text") or "")
        left, top, width, height = ctx.to_inches(comp.x, comp.y, comp.w, comp.h)
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        p.alignment = _align(comp.style.align or "left")
        font = run.font
        font.name = comp.style.font_family or ctx.theme.typography.font_family
        font.size = Pt(comp.style.font_size or ctx.theme.typography.title_pt)
        font.bold = True if comp.style.bold is None else comp.style.bold
        font.italic = bool(comp.style.italic) if comp.style.italic is not None else False
        font.color.rgb = _hex_to_rgb(comp.style.color or ctx.theme.colors.text)


class SubtitleRenderer:
    type = "Subtitle"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        text = str(comp.props.get("text") or "")
        left, top, width, height = ctx.to_inches(comp.x, comp.y, comp.w, comp.h)
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        p.alignment = _align(comp.style.align or "left")
        font = run.font
        font.name = comp.style.font_family or ctx.theme.typography.font_family
        font.size = Pt(comp.style.font_size or ctx.theme.typography.subtitle_pt)
        font.bold = bool(comp.style.bold) if comp.style.bold is not None else False
        font.italic = bool(comp.style.italic) if comp.style.italic is not None else False
        font.color.rgb = _hex_to_rgb(comp.style.color or ctx.theme.colors.muted)


class TextRenderer:
    type = "Text"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        text = str(comp.props.get("text") or "")
        left, top, width, height = ctx.to_inches(comp.x, comp.y, comp.w, comp.h)
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, line in enumerate(text.splitlines() or [""]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.alignment = _align(comp.style.align or "left")
            p.font.name = comp.style.font_family or ctx.theme.typography.font_family
            p.font.size = Pt(comp.style.font_size or ctx.theme.typography.body_pt)
            p.font.color.rgb = _hex_to_rgb(comp.style.color or ctx.theme.colors.text)


class BulletListRenderer:
    type = "BulletList"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        items = comp.props.get("items") or []
        if isinstance(items, str):
            items = [items]
        left, top, width, height = ctx.to_inches(comp.x, comp.y, comp.w, comp.h)
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(item)
            p.level = 0
            p.font.name = comp.style.font_family or ctx.theme.typography.font_family
            p.font.size = Pt(comp.style.font_size or ctx.theme.typography.body_pt)
            p.font.color.rgb = _hex_to_rgb(comp.style.color or ctx.theme.colors.text)


class DividerRenderer:
    type = "Divider"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        title = str(comp.props.get("title") or "")
        subtitle = comp.props.get("subtitle")
        left, top, width, height = ctx.to_inches(comp.x, comp.y, comp.w, comp.h)

        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = _hex_to_rgb(ctx.theme.colors.background)
        bg.line.fill.background()

        t_left, t_top, t_w, t_h = ctx.to_inches(comp.x, comp.y + comp.h * 0.25, comp.w, comp.h * 0.4)
        title_box = slide.shapes.add_textbox(t_left, t_top, t_w, t_h)
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = ctx.theme.typography.font_family
        p.font.size = Pt(ctx.theme.typography.title_pt)
        p.font.bold = True
        p.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)

        if subtitle:
            s_left, s_top, s_w, s_h = ctx.to_inches(comp.x, comp.y + comp.h * 0.60, comp.w, comp.h * 0.2)
            sub_box = slide.shapes.add_textbox(s_left, s_top, s_w, s_h)
            tf2 = sub_box.text_frame
            tf2.clear()
            p2 = tf2.paragraphs[0]
            p2.text = str(subtitle)
            p2.alignment = PP_ALIGN.CENTER
            p2.font.name = ctx.theme.typography.font_family
            p2.font.size = Pt(ctx.theme.typography.subtitle_pt)
            p2.font.color.rgb = _hex_to_rgb(ctx.theme.colors.muted)


class QuoteRenderer:
    type = "Quote"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        quote = str(comp.props.get("quote") or "")
        author = comp.props.get("author")
        left, top, width, height = ctx.to_inches(comp.x, comp.y, comp.w, comp.h)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = _hex_to_rgb(comp.style.background or ctx.theme.colors.surface)
        card.line.color.rgb = _hex_to_rgb(comp.style.border_color or ctx.theme.colors.primary)
        card.line.width = Pt(comp.style.border_width or 2)

        tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), width - Inches(0.6), height - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"“{quote}”"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = comp.style.font_family or ctx.theme.typography.font_family
        p.font.size = Pt(comp.style.font_size or ctx.theme.typography.subtitle_pt)
        p.font.italic = True
        p.font.color.rgb = _hex_to_rgb(comp.style.color or ctx.theme.colors.text)
        if author:
            p2 = tf.add_paragraph()
            p2.text = f"— {author}"
            p2.alignment = PP_ALIGN.CENTER
            p2.font.name = comp.style.font_family or ctx.theme.typography.font_family
            p2.font.size = Pt(ctx.theme.typography.body_pt)
            p2.font.color.rgb = _hex_to_rgb(ctx.theme.colors.muted)


class KpiCardsRenderer:
    type = "KpiCards"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        items = comp.props.get("items") or []
        if not isinstance(items, list):
            return
        cols = min(3, max(1, len(items)))
        rows = (len(items) + cols - 1) // cols
        gap_px = ctx.theme.spacing.gap_px
        card_w = (comp.w - gap_px * (cols - 1)) / cols
        card_h = (comp.h - gap_px * (rows - 1)) / rows if rows else comp.h
        for idx, item in enumerate(items):
            r = idx // cols
            c = idx % cols
            x = comp.x + c * (card_w + gap_px)
            y = comp.y + r * (card_h + gap_px)
            left, top, width, height = ctx.to_inches(x, y, card_w, card_h)
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
            card.fill.solid()
            card.fill.fore_color.rgb = _hex_to_rgb(comp.style.background or ctx.theme.colors.surface)
            card.line.color.rgb = _hex_to_rgb(comp.style.border_color or ctx.theme.colors.border)
            card.line.width = Pt(comp.style.border_width or 1)

            label = str((item or {}).get("label") or "")
            value = str((item or {}).get("value") or "")
            unit = (item or {}).get("unit")
            delta = (item or {}).get("delta")
            value_text = value + (f" {unit}" if unit else "")
            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
            tf = tb.text_frame
            tf.clear()
            p1 = tf.paragraphs[0]
            p1.text = label
            p1.font.name = ctx.theme.typography.font_family
            p1.font.size = Pt(ctx.theme.typography.small_pt)
            p1.font.color.rgb = _hex_to_rgb(ctx.theme.colors.muted)
            p2 = tf.add_paragraph()
            p2.text = value_text
            p2.font.name = ctx.theme.typography.font_family
            p2.font.size = Pt(ctx.theme.typography.title_pt)
            p2.font.bold = True
            p2.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)
            if delta:
                p3 = tf.add_paragraph()
                p3.text = str(delta)
                p3.font.name = ctx.theme.typography.font_family
                p3.font.size = Pt(ctx.theme.typography.small_pt)
                p3.font.color.rgb = _hex_to_rgb(ctx.theme.colors.primary)


class TimelineRenderer:
    type = "Timeline"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        events = comp.props.get("events") or []
        if not isinstance(events, list):
            return
        if not events:
            return
        left, top, width, height = ctx.to_inches(comp.x, comp.y, comp.w, comp.h)
        y_mid = comp.y + comp.h * 0.5
        line_left, line_top, line_w, line_h = ctx.to_inches(comp.x, y_mid, comp.w, 2)
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, line_left, line_top, line_w, line_h)
        line.fill.solid()
        line.fill.fore_color.rgb = _hex_to_rgb(ctx.theme.colors.border)
        line.line.fill.background()

        n = len(events)
        for i, e in enumerate(events):
            x = comp.x + (i + 0.5) * (comp.w / n)
            dot_left, dot_top, dot_w, dot_h = ctx.to_inches(x - 6, y_mid - 6, 12, 12)
            dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, dot_left, dot_top, dot_w, dot_h)
            dot.fill.solid()
            dot.fill.fore_color.rgb = _hex_to_rgb(ctx.theme.colors.primary)
            dot.line.fill.background()
            label = str((e or {}).get("label") or "")
            date = (e or {}).get("date")
            detail = (e or {}).get("detail")
            text = label
            if date:
                text = f"{date} · {text}"
            if detail:
                text = f"{text}\n{detail}"
            tb_left, tb_top, tb_w, tb_h = ctx.to_inches(x - (comp.w / n) * 0.45, comp.y, (comp.w / n) * 0.9, comp.h * 0.45)
            tb = slide.shapes.add_textbox(tb_left, tb_top, tb_w, tb_h)
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.CENTER
            p.font.name = ctx.theme.typography.font_family
            p.font.size = Pt(ctx.theme.typography.small_pt)
            p.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)


class ComparisonTableRenderer:
    type = "ComparisonTable"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        left_data = comp.props.get("left") or {}
        right_data = comp.props.get("right") or {}
        gap_px = ctx.theme.spacing.gap_px
        col_w = (comp.w - gap_px) / 2
        boxes = [
            (comp.x, comp.y, col_w, comp.h, left_data),
            (comp.x + col_w + gap_px, comp.y, col_w, comp.h, right_data),
        ]
        for x, y, w, h, data in boxes:
            l, t, ww, hh = ctx.to_inches(x, y, w, h)
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, ww, hh)
            card.fill.solid()
            card.fill.fore_color.rgb = _hex_to_rgb(ctx.theme.colors.surface)
            card.line.color.rgb = _hex_to_rgb(ctx.theme.colors.border)
            card.line.width = Pt(1)
            title = str((data or {}).get("title") or "")
            bullets = (data or {}).get("bullets") or []
            tb = slide.shapes.add_textbox(l + Inches(0.2), t + Inches(0.2), ww - Inches(0.4), hh - Inches(0.4))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = title
            p0.font.name = ctx.theme.typography.font_family
            p0.font.size = Pt(ctx.theme.typography.subtitle_pt)
            p0.font.bold = True
            p0.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)
            for b in bullets:
                p = tf.add_paragraph()
                p.text = str(b)
                p.level = 0
                p.font.name = ctx.theme.typography.font_family
                p.font.size = Pt(ctx.theme.typography.body_pt)
                p.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)


class SwotRenderer:
    type = "Swot"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        gap_px = ctx.theme.spacing.gap_px
        cell_w = (comp.w - gap_px) / 2
        cell_h = (comp.h - gap_px) / 2
        cells = [
            ("Strengths", comp.props.get("strengths") or [], comp.x, comp.y),
            ("Weaknesses", comp.props.get("weaknesses") or [], comp.x + cell_w + gap_px, comp.y),
            ("Opportunities", comp.props.get("opportunities") or [], comp.x, comp.y + cell_h + gap_px),
            ("Threats", comp.props.get("threats") or [], comp.x + cell_w + gap_px, comp.y + cell_h + gap_px),
        ]
        for title, items, x, y in cells:
            l, t, w, h = ctx.to_inches(x, y, cell_w, cell_h)
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
            card.fill.solid()
            card.fill.fore_color.rgb = _hex_to_rgb(ctx.theme.colors.surface)
            card.line.color.rgb = _hex_to_rgb(ctx.theme.colors.border)
            card.line.width = Pt(1)
            tb = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.15), w - Inches(0.3), h - Inches(0.3))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = title
            p0.font.name = ctx.theme.typography.font_family
            p0.font.size = Pt(ctx.theme.typography.subtitle_pt)
            p0.font.bold = True
            p0.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)
            for it in items:
                p = tf.add_paragraph()
                p.text = str(it)
                p.level = 0
                p.font.name = ctx.theme.typography.font_family
                p.font.size = Pt(ctx.theme.typography.small_pt)
                p.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)


class RoadmapRenderer:
    type = "Roadmap"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        phases = comp.props.get("phases") or []
        if not isinstance(phases, list) or not phases:
            return
        gap_px = ctx.theme.spacing.gap_px
        n = len(phases)
        box_w = (comp.w - gap_px * (n - 1)) / n
        for i, ph in enumerate(phases):
            x = comp.x + i * (box_w + gap_px)
            l, t, w, h = ctx.to_inches(x, comp.y, box_w, comp.h)
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
            card.fill.solid()
            card.fill.fore_color.rgb = _hex_to_rgb(ctx.theme.colors.surface)
            card.line.color.rgb = _hex_to_rgb(ctx.theme.colors.border)
            card.line.width = Pt(1)
            name = str((ph or {}).get("name") or "")
            timeframe = (ph or {}).get("timeframe")
            deliverables = (ph or {}).get("deliverables") or []
            tb = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.15), w - Inches(0.3), h - Inches(0.3))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = name
            p0.font.name = ctx.theme.typography.font_family
            p0.font.size = Pt(ctx.theme.typography.subtitle_pt)
            p0.font.bold = True
            p0.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)
            if timeframe:
                p1 = tf.add_paragraph()
                p1.text = str(timeframe)
                p1.font.name = ctx.theme.typography.font_family
                p1.font.size = Pt(ctx.theme.typography.small_pt)
                p1.font.color.rgb = _hex_to_rgb(ctx.theme.colors.muted)
            for d in deliverables:
                p = tf.add_paragraph()
                p.text = str(d)
                p.level = 0
                p.font.name = ctx.theme.typography.font_family
                p.font.size = Pt(ctx.theme.typography.small_pt)
                p.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)


class ProcessFlowRenderer:
    type = "ProcessFlow"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        steps = comp.props.get("steps") or []
        if not isinstance(steps, list) or not steps:
            return
        gap_px = ctx.theme.spacing.gap_px
        n = len(steps)
        box_w = (comp.w - gap_px * (n - 1)) / n
        for i, st in enumerate(steps):
            x = comp.x + i * (box_w + gap_px)
            l, t, w, h = ctx.to_inches(x, comp.y, box_w, comp.h)
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
            card.fill.solid()
            card.fill.fore_color.rgb = _hex_to_rgb(ctx.theme.colors.surface)
            card.line.color.rgb = _hex_to_rgb(ctx.theme.colors.border)
            card.line.width = Pt(1)
            name = str((st or {}).get("name") or "")
            detail = (st or {}).get("detail")
            tb = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.15), w - Inches(0.3), h - Inches(0.3))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = name
            p0.alignment = PP_ALIGN.CENTER
            p0.font.name = ctx.theme.typography.font_family
            p0.font.size = Pt(ctx.theme.typography.subtitle_pt)
            p0.font.bold = True
            p0.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)
            if detail:
                p1 = tf.add_paragraph()
                p1.text = str(detail)
                p1.alignment = PP_ALIGN.CENTER
                p1.font.name = ctx.theme.typography.font_family
                p1.font.size = Pt(ctx.theme.typography.small_pt)
                p1.font.color.rgb = _hex_to_rgb(ctx.theme.colors.muted)


class ChartRenderer:
    type = "Chart"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        chart_type = str(comp.props.get("chartType") or comp.props.get("chart_type") or "bar")
        labels = comp.props.get("labels") or []
        series = comp.props.get("series") or []
        if not labels or not series:
            return
        chart_data = CategoryChartData()
        chart_data.categories = [str(x) for x in labels]
        for s in series:
            name = str((s or {}).get("name") or "Series")
            values = (s or {}).get("values") or []
            chart_data.add_series(name, values)
        left, top, width, height = ctx.to_inches(comp.x, comp.y, comp.w, comp.h)
        xl_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        if chart_type == "line":
            xl_type = XL_CHART_TYPE.LINE
        if chart_type == "pie":
            xl_type = XL_CHART_TYPE.PIE
        chart = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data).chart
        chart.has_legend = True
        chart.legend.include_in_layout = False


class MultiColumnRenderer:
    type = "MultiColumn"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        columns = comp.props.get("columns") or []
        if not isinstance(columns, list) or not columns:
            return
        gap_px = ctx.theme.spacing.gap_px
        n = len(columns)
        box_w = (comp.w - gap_px * (n - 1)) / n
        for i, col in enumerate(columns):
            x = comp.x + i * (box_w + gap_px)
            l, t, w, h = ctx.to_inches(x, comp.y, box_w, comp.h)
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
            card.fill.solid()
            card.fill.fore_color.rgb = _hex_to_rgb(ctx.theme.colors.surface)
            card.line.color.rgb = _hex_to_rgb(ctx.theme.colors.border)
            card.line.width = Pt(1)
            title = str((col or {}).get("title") or "")
            bullets = (col or {}).get("bullets") or []
            tb = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.15), w - Inches(0.3), h - Inches(0.3))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = title
            p0.font.name = ctx.theme.typography.font_family
            p0.font.size = Pt(ctx.theme.typography.subtitle_pt)
            p0.font.bold = True
            p0.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)
            for b in bullets:
                p = tf.add_paragraph()
                p.text = str(b)
                p.level = 0
                p.font.name = ctx.theme.typography.font_family
                p.font.size = Pt(ctx.theme.typography.small_pt)
                p.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)


class TeamCardsRenderer:
    type = "TeamCards"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        members = comp.props.get("members") or []
        if not isinstance(members, list) or not members:
            return
        cols = min(3, max(1, len(members)))
        rows = (len(members) + cols - 1) // cols
        gap_px = ctx.theme.spacing.gap_px
        card_w = (comp.w - gap_px * (cols - 1)) / cols
        card_h = (comp.h - gap_px * (rows - 1)) / rows if rows else comp.h
        for idx, m in enumerate(members):
            r = idx // cols
            c = idx % cols
            x = comp.x + c * (card_w + gap_px)
            y = comp.y + r * (card_h + gap_px)
            l, t, w, h = ctx.to_inches(x, y, card_w, card_h)
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
            card.fill.solid()
            card.fill.fore_color.rgb = _hex_to_rgb(ctx.theme.colors.surface)
            card.line.color.rgb = _hex_to_rgb(ctx.theme.colors.border)
            card.line.width = Pt(1)
            name = str((m or {}).get("name") or "")
            role = (m or {}).get("role")
            highlights = (m or {}).get("highlights") or []
            tb = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.15), w - Inches(0.3), h - Inches(0.3))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = name
            p0.font.name = ctx.theme.typography.font_family
            p0.font.size = Pt(ctx.theme.typography.subtitle_pt)
            p0.font.bold = True
            p0.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)
            if role:
                p1 = tf.add_paragraph()
                p1.text = str(role)
                p1.font.name = ctx.theme.typography.font_family
                p1.font.size = Pt(ctx.theme.typography.small_pt)
                p1.font.color.rgb = _hex_to_rgb(ctx.theme.colors.muted)
            for hlt in highlights:
                p = tf.add_paragraph()
                p.text = str(hlt)
                p.level = 0
                p.font.name = ctx.theme.typography.font_family
                p.font.size = Pt(ctx.theme.typography.small_pt)
                p.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)


class ArchitectureDiagramRenderer:
    type = "ArchitectureDiagram"

    def render(self, slide, comp: ComponentSpec, ctx: PptxRenderContext) -> None:
        layers = comp.props.get("layers") or []
        if not isinstance(layers, list) or not layers:
            return
        gap_px = ctx.theme.spacing.gap_px
        n = len(layers)
        box_h = (comp.h - gap_px * (n - 1)) / n
        for i, layer in enumerate(layers):
            y = comp.y + i * (box_h + gap_px)
            l, t, w, h = ctx.to_inches(comp.x, y, comp.w, box_h)
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
            card.fill.solid()
            card.fill.fore_color.rgb = _hex_to_rgb(ctx.theme.colors.surface)
            card.line.color.rgb = _hex_to_rgb(ctx.theme.colors.border)
            card.line.width = Pt(1)
            name = str((layer or {}).get("name") or "")
            items = (layer or {}).get("items") or []
            tb = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.10), w - Inches(0.3), h - Inches(0.2))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.text = name
            p0.font.name = ctx.theme.typography.font_family
            p0.font.size = Pt(ctx.theme.typography.subtitle_pt)
            p0.font.bold = True
            p0.font.color.rgb = _hex_to_rgb(ctx.theme.colors.text)
            if items:
                p1 = tf.add_paragraph()
                p1.text = " · ".join([str(x) for x in items])
                p1.font.name = ctx.theme.typography.font_family
                p1.font.size = Pt(ctx.theme.typography.small_pt)
                p1.font.color.rgb = _hex_to_rgb(ctx.theme.colors.muted)
