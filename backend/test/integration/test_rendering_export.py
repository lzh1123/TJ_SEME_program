from __future__ import annotations

from pptx import Presentation

from ppt_backend.domain.dsl import PresentationDSL
from ppt_backend.domain.theme import get_theme_tokens
from ppt_backend.exporters.pptx_components import build_component_renderer_registry
from ppt_backend.exporters.pptx_exporter import PptxExporter
from ppt_backend.services.rendering.compiler import RenderCompiler
from ppt_backend.services.rendering.registry import (
    build_layout_registry,
    build_slide_composer_registry,
)

from backend.test.sample_deck_fixtures import sample_full_dsl


def test_dsl_render_tree_and_pptx_export_round_trip(tmp_path):
    dsl = PresentationDSL.model_validate(sample_full_dsl())
    compiler = RenderCompiler(build_slide_composer_registry(), build_layout_registry())
    tree = compiler.compile("pres_validation", dsl, get_theme_tokens(dsl.theme))

    assert len(tree.slides) == len(dsl.slides)
    assert {slide.intent for slide in dsl.slides} == {
        "cover",
        "agenda",
        "text",
        "timeline",
        "kpi",
        "comparison",
        "swot",
        "roadmap",
        "process_flow",
        "chart",
        "multi_column",
        "architecture",
        "quote",
        "divider",
        "team",
    }
    for slide in tree.slides:
        assert slide.components
        for component in slide.components:
            assert 0 <= component.x <= slide.width
            assert 0 <= component.y <= slide.height
            assert component.w > 0
            assert component.h > 0
            assert component.x + component.w <= slide.width + 0.01
            assert component.y + component.h <= slide.height + 0.01

    out = PptxExporter(build_component_renderer_registry()).export(tree, tmp_path / "validation_export.pptx")
    assert out.exists() and out.stat().st_size > 0
    prs = Presentation(str(out))
    assert len(prs.slides) == len(dsl.slides)
