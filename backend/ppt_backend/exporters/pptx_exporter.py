from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ..domain.render_tree import RenderTree
from ..registry.base import Registry
from .pptx_components import PptxRenderContext, ComponentRenderer


class PptxExporter:
    def __init__(self, component_renderers: Registry[ComponentRenderer]):
        self._renderers = component_renderers

    def export(self, tree: RenderTree, out_path: Path) -> Path:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        theme = tree.theme_tokens
        ctx = PptxRenderContext(theme=theme)

        blank_layout = prs.slide_layouts[6]
        for s in tree.slides:
            slide = prs.slides.add_slide(blank_layout)
            for comp in sorted(s.components, key=lambda c: (c.z, c.id)):
                renderer = self._renderers.get(comp.type)
                renderer.render(slide, comp, ctx)
            if s.notes:
                notes_text = "\n".join([str(n) for n in s.notes if n is not None])
                slide.notes_slide.notes_text_frame.text = notes_text

        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        return out_path
