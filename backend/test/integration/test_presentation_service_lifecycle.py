from __future__ import annotations

from pptx import Presentation

from ppt_backend.domain.render_tree import ComponentPatch

from backend.test.sample_deck_fixtures import make_presentation_service


def test_presentation_service_editing_theme_reorder_and_export_lifecycle(tmp_path):
    service = make_presentation_service(tmp_path)
    bundle = service.create_from_outline(
        topic="Presentation lifecycle",
        outline={
            "title": "Presentation lifecycle",
            "slides": [
                {"intent": "cover", "title": "Lifecycle", "subtitle": "Start", "notes": "intro"},
                {"intent": "text", "title": "Content", "content": "Generated content", "bullets": ["A", "B"]},
                {"intent": "agenda", "title": "Agenda", "items": ["One", "Two"]},
            ],
        },
        theme="paper_light",
    )

    component_id = bundle.render_tree.slides[1].components[0].id
    patched = service.patch_component(
        bundle.meta.id,
        component_id,
        ComponentPatch(props={"text": "Updated title"}, style={"fontSize": 36}),
    )
    updated_component = next(c for c in patched.render_tree.slides[1].components if c.id == component_id)
    assert updated_component.props["text"] == "Updated title"
    assert updated_component.style.font_size == 36

    reversed_ids = [slide.id for slide in reversed(patched.render_tree.slides)]
    reordered = service.reorder_slides(bundle.meta.id, reversed_ids)
    assert [slide.id for slide in reordered.render_tree.slides] == reversed_ids

    themed = service.switch_theme(bundle.meta.id, "modern_blue")
    assert themed.dsl.theme == "modern_blue"
    assert themed.render_tree.theme_name == "modern_blue"

    pptx_path = service.export_pptx(bundle.meta.id)
    assert pptx_path.exists()
    assert len(Presentation(str(pptx_path)).slides) == 3
