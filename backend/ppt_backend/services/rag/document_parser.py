from __future__ import annotations

from pathlib import Path


SUPPORTED_DOCUMENT_SUFFIXES = {".pdf", ".docx", ".txt", ".md", ".pptx"}


class DocumentParseError(ValueError):
    pass


def parse_document(path: Path, suffix: str | None = None) -> str:
    suffix = (suffix or path.suffix).lower()
    if suffix in {".txt", ".md"}:
        return _read_text(path)
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    raise DocumentParseError(
        f"Unsupported file type: {suffix}. Supported: {', '.join(sorted(SUPPORTED_DOCUMENT_SUFFIXES))}"
    )


def _read_text(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="ignore")


def _read_pdf(path: Path) -> str:
    parts: list[str] = []
    try:
        import fitz

        doc = fitz.open(str(path))
        try:
            for page in doc:
                text = page.get_text()
                if text and text.strip():
                    parts.append(text.strip())
        finally:
            doc.close()
    except ImportError:
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            for page in reader.pages:
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(text.strip())
        except ImportError as exc:
            raise DocumentParseError("PDF parsing requires PyMuPDF or pypdf") from exc

    return "\n\n".join(parts)


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise DocumentParseError("DOCX parsing requires python-docx") from exc

    doc = Document(str(path))
    parts: list[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentParseError("PPTX parsing requires python-pptx") from exc

    prs = Presentation(str(path))
    parts: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                slide_parts.append(text.strip())
        if slide_parts:
            parts.append(f"Slide {index}\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def compact_document_text(content: str, max_chars: int = 12000) -> str:
    content = "\n".join(line.rstrip() for line in content.splitlines()).strip()
    if len(content) <= max_chars:
        return content
    head = content[: int(max_chars * 0.55)]
    middle_start = max(0, len(content) // 2 - int(max_chars * 0.1))
    middle = content[middle_start : middle_start + int(max_chars * 0.2)]
    tail = content[-int(max_chars * 0.25) :]
    return f"{head}\n\n...[content truncated]...\n\n{middle}\n\n...[content truncated]...\n\n{tail}"
